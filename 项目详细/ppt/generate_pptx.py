"""生成项目报告PPT"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 配置
IKB_BLUE = RGBColor(0x00, 0x2F, 0xA7)
DARK_BG = RGBColor(0x0A, 0x0A, 0x0A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xFA, 0xFA, 0xF8)
GREY = RGBColor(0x73, 0x73, 0x73)
BORDER_GREY = RGBColor(0xE0, 0xE0, 0xE0)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

IMG_DIR = "images"


def add_bg(slide, color=DARK_BG):
    """设置幻灯片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False,
                 color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    if font_name:
        p.font.name = font_name
    return txBox


def add_image(slide, img_path, left, top, width, height=None):
    """添加图片"""
    if os.path.exists(img_path):
        if height:
            slide.shapes.add_picture(img_path, Inches(left), Inches(top), Inches(width), Inches(height))
        else:
            slide.shapes.add_picture(img_path, Inches(left), Inches(top), Inches(width))


def add_shape_bg(slide, left, top, width, height, color):
    """添加矩形背景"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ==================== Slide 1: 封面 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
add_bg(slide, DARK_BG)

add_text_box(slide, 1, 0.5, 4, 0.4, "项目报告 · 2026.07", font_size=12, color=GREY)
add_text_box(slide, 1, 1.5, 11, 0.5, "全栈开发 · AI智能客服", font_size=14, color=IKB_BLUE)
add_text_box(slide, 1, 2.2, 11, 1.5, "公司桌面IT服务台", font_size=60, bold=True, color=WHITE)
add_text_box(slide, 1, 4.0, 11, 0.8, "AI智能客服 + RAG知识库 + 三层会话记忆", font_size=28, color=RGBColor(0xBB, 0xBB, 0xBB))
add_text_box(slide, 1, 5.5, 11, 0.4, "FastAPI  ·  Vue 3  ·  ChromaDB  ·  WebSocket  ·  Docker", font_size=14, color=GREY)

# 底部线条
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(6.8), Inches(11.3), Pt(1))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x33)
shape.line.fill.background()
add_text_box(slide, 1, 6.9, 5, 0.3, "公司桌面IT服务台 · AI智能版", font_size=10, color=GREY)
add_text_box(slide, 10, 6.9, 2, 0.3, "01 / 12", font_size=10, color=GREY, alignment=PP_ALIGN.RIGHT)


# ==================== Slide 2: 项目背景 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GREY)

add_text_box(slide, 0.8, 0.5, 3, 0.3, "01  ·  项目背景", font_size=11, color=GREY)
add_text_box(slide, 0.8, 1.2, 5, 0.3, "业务痛点", font_size=14, color=IKB_BLUE)
add_text_box(slide, 0.8, 1.7, 11, 0.8, "为什么需要AI智能客服？", font_size=44, bold=True, color=DARK_BG)

# 三个痛点卡片
cards = [
    ("响应延迟", "用户提交工单后需等待客服响应，首次回复延迟平均15分钟+，高峰期更长"),
    ("重复劳动", "常见问题（蓝屏、密码重置、打印机故障）重复率高达60%+，客服重复解答"),
    ("知识流失", "已解决的问题无法复用，同类问题反复处理，缺乏知识库沉淀机制"),
]
for i, (title, desc) in enumerate(cards):
    x = 0.8 + i * 4.0
    add_shape_bg(slide, x, 3.2, 3.6, 2.8, WHITE)
    add_shape_bg(slide, x, 3.2, 3.6, 0.06, IKB_BLUE)
    add_text_box(slide, x + 0.3, 3.5, 3, 0.4, title, font_size=18, bold=True, color=IKB_BLUE)
    add_text_box(slide, x + 0.3, 4.2, 3, 1.5, desc, font_size=14, color=RGBColor(0x52, 0x52, 0x52))

add_text_box(slide, 0.8, 6.9, 5, 0.3, "项目背景", font_size=10, color=GREY)
add_text_box(slide, 10, 6.9, 2, 0.3, "02 / 12", font_size=10, color=GREY, alignment=PP_ALIGN.RIGHT)


# ==================== Slide 3: 系统架构 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_text_box(slide, 0.8, 0.5, 3, 0.3, "02  ·  系统架构", font_size=11, color=GREY)
add_text_box(slide, 0.8, 1.2, 5, 0.3, "整体设计", font_size=14, color=IKB_BLUE)
add_text_box(slide, 0.8, 1.7, 11, 0.8, "四端分离 + AI中台", font_size=44, bold=True, color=WHITE)

# 架构层
layers = [
    ("用户端 :5173", "服务台 · AI客服 · 工单 · 聊天室 · 评价"),
    ("客服端 :5174", "工作台 · 工单池 · 实时聊天 · SLA管理 · 转派"),
    ("统计端 :5176", "数据概览 · 趋势分析 · 绩效排名 · 报表导出"),
    ("管理端 :5175", "用户管理 · 权限控制 · 分类配置 · 审计日志"),
]
for i, (title, desc) in enumerate(layers):
    y = 3.0 + i * 0.85
    add_shape_bg(slide, 0.8, y, 11.7, 0.7, RGBColor(0x1A, 0x1A, 0x1A))
    add_text_box(slide, 1.0, y + 0.1, 2, 0.4, title, font_size=14, bold=True, color=WHITE)
    add_text_box(slide, 3.5, y + 0.1, 9, 0.4, desc, font_size=14, color=RGBColor(0xBB, 0xBB, 0xBB))

# AI中台
add_shape_bg(slide, 0.8, 6.5, 11.7, 0.7, RGBColor(0x00, 0x1A, 0x4A))
add_text_box(slide, 1.0, 6.6, 2, 0.4, "AI中台 :8000", font_size=14, bold=True, color=IKB_BLUE)
add_text_box(slide, 3.5, 6.6, 9, 0.4, "RAG管道 · 向量检索 · LLM生成 · 会话记忆 · 知识库", font_size=14, color=RGBColor(0xBB, 0xBB, 0xBB))

add_text_box(slide, 0.8, 7.1, 5, 0.3, "系统架构", font_size=10, color=GREY)
add_text_box(slide, 10, 7.1, 2, 0.3, "03 / 12", font_size=10, color=GREY, alignment=PP_ALIGN.RIGHT)


# ==================== Slide 4: AI核心能力 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GREY)

add_text_box(slide, 0.8, 0.5, 3, 0.3, "03  ·  AI核心能力", font_size=11, color=GREY)
add_text_box(slide, 0.8, 1.2, 5, 0.3, "RAG检索增强生成", font_size=14, color=IKB_BLUE)
add_text_box(slide, 0.8, 1.7, 5, 0.8, "AI智能客服", font_size=44, bold=True, color=DARK_BG)
add_text_box(slide, 0.8, 2.8, 5, 1.2, "基于已解决工单和FAQ文档知识库，智能检索相关解决方案，自动生成专业回答", font_size=18, color=RGBColor(0x52, 0x52, 0x52))

# 数据卡片
add_shape_bg(slide, 0.8, 4.5, 2.5, 1.8, WHITE)
add_shape_bg(slide, 0.8, 4.5, 2.5, 0.06, IKB_BLUE)
add_text_box(slide, 1.1, 4.8, 2, 0.3, "知识库文档", font_size=12, color=GREY)
add_text_box(slide, 1.1, 5.2, 2, 0.6, "100+", font_size=48, bold=True, color=IKB_BLUE)
add_text_box(slide, 1.1, 5.9, 2, 0.3, "常见问题FAQ", font_size=12, color=GREY)

add_shape_bg(slide, 3.6, 4.5, 2.5, 1.8, WHITE)
add_shape_bg(slide, 3.6, 4.5, 2.5, 0.06, IKB_BLUE)
add_text_box(slide, 3.9, 4.8, 2, 0.3, "检索精度", font_size=12, color=GREY)
add_text_box(slide, 3.9, 5.2, 2, 0.6, "0.5", font_size=48, bold=True, color=IKB_BLUE)
add_text_box(slide, 3.9, 5.9, 2, 0.3, "相似度阈值", font_size=12, color=GREY)

# 截图
add_image(slide, "images/03-chat.png", 7, 1.5, 5.5, 5.0)

add_text_box(slide, 0.8, 6.9, 5, 0.3, "AI核心能力", font_size=10, color=GREY)
add_text_box(slide, 10, 6.9, 2, 0.3, "04 / 12", font_size=10, color=GREY, alignment=PP_ALIGN.RIGHT)


# ==================== Slide 5: 三层记忆系统 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(0xF0, 0xF0, 0xEE))

add_text_box(slide, 0.8, 0.5, 3, 0.3, "04  ·  记忆系统", font_size=11, color=GREY)
add_text_box(slide, 0.8, 1.2, 5, 0.3, "三层会话记忆架构", font_size=14, color=IKB_BLUE)
add_text_box(slide, 0.8, 1.7, 11, 0.8, "精准对话，永不失忆", font_size=44, bold=True, color=DARK_BG)

# 三层记忆卡片
layers = [
    ("滑动窗口", "5轮", "最近5轮原始对话精确回放，确保上下文连贯"),
    ("会话摘要", "AUTO", "旧对话LLM自动浓缩为摘要，防止长对话信息丢失"),
    ("会话元数据", "4D", "设备型号/操作系统/问题分类/场景，每3轮自动提取"),
]
for i, (title, value, desc) in enumerate(layers):
    x = 0.8 + i * 4.0
    add_shape_bg(slide, x, 3.2, 3.6, 3.0, WHITE)
    add_shape_bg(slide, x, 3.2, 3.6, 0.08, IKB_BLUE)
    add_text_box(slide, x + 0.3, 3.5, 3, 0.4, title, font_size=14, color=IKB_BLUE)
    add_text_box(slide, x + 0.3, 4.2, 3, 0.8, value, font_size=48, bold=True, color=DARK_BG)
    add_text_box(slide, x + 0.3, 5.3, 3, 0.8, desc, font_size=14, color=RGBColor(0x52, 0x52, 0x52))

add_text_box(slide, 0.8, 6.5, 11, 0.3, "Redis持久化 · TTL 30分钟 · 内存fallback双写", font_size=12, color=GREY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 0.8, 7.1, 5, 0.3, "记忆系统", font_size=10, color=GREY)
add_text_box(slide, 10, 7.1, 2, 0.3, "05 / 12", font_size=10, color=GREY, alignment=PP_ALIGN.RIGHT)


# ==================== Slide 6: 技术栈 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GREY)

add_text_box(slide, 0.8, 0.5, 3, 0.3, "05  ·  技术栈", font_size=11, color=GREY)
add_text_box(slide, 0.8, 1.2, 5, 0.3, "核心技术选型", font_size=14, color=IKB_BLUE)
add_text_box(slide, 0.8, 1.7, 11, 0.8, "全栈技术方案", font_size=44, bold=True, color=DARK_BG)

# 四个技术卡片
techs = [
    ("后端", "FastAPI + async SQLAlchemy", "异步架构 · SQLite/MySQL · JWT认证 · APScheduler"),
    ("前端", "Vue 3 + Element Plus", "Pinia状态管理 · ECharts图表 · WebSocket实时通信"),
    ("AI", "RAG + ChromaDB + LLM", "BGE嵌入 · 向量检索 · DeepSeek/Qwen · SSE流式"),
    ("DevOps", "Docker + GitHub Actions", "CI/CD · 73个测试 · Redis缓存 · 健康检查"),
]
for i, (title, main, sub) in enumerate(techs):
    x = 0.8 + i * 3.05
    add_shape_bg(slide, x, 3.2, 2.8, 3.0, WHITE)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(3.2), Inches(2.8), Pt(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = IKB_BLUE
    shape.line.fill.background()
    add_text_box(slide, x + 0.2, 3.5, 2.4, 0.3, title, font_size=14, color=IKB_BLUE)
    add_text_box(slide, x + 0.2, 4.0, 2.4, 0.5, main, font_size=16, bold=True, color=DARK_BG)
    add_text_box(slide, x + 0.2, 4.8, 2.4, 1.2, sub, font_size=12, color=RGBColor(0x73, 0x73, 0x73))

add_text_box(slide, 0.8, 7.1, 5, 0.3, "技术栈", font_size=10, color=GREY)
add_text_box(slide, 10, 7.1, 2, 0.3, "06 / 12", font_size=10, color=GREY, alignment=PP_ALIGN.RIGHT)


# ==================== Slide 7: 用户服务台 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_text_box(slide, 0.8, 0.5, 3, 0.3, "06  ·  用户服务台", font_size=11, color=GREY)
add_text_box(slide, 0.8, 1.0, 5, 0.3, "用户端 · :5173", font_size=14, color=IKB_BLUE)
add_text_box(slide, 0.8, 1.5, 5, 0.8, "AI智能客服 + 一键转人工", font_size=36, bold=True, color=WHITE)

# 截图
add_image(slide, "images/01-cover.jpg", 0.8, 2.8, 11.7, 4.2)

add_text_box(slide, 0.8, 7.1, 5, 0.3, "用户服务台", font_size=10, color=GREY)
add_text_box(slide, 10, 7.1, 2, 0.3, "07 / 12", font_size=10, color=GREY, alignment=PP_ALIGN.RIGHT)


# ==================== Slide 8: ITSM客服端 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GREY)

add_text_box(slide, 0.8, 0.5, 3, 0.3, "07  ·  ITSM客服端", font_size=11, color=GREY)
add_text_box(slide, 0.8, 1.2, 5, 0.3, "工单全生命周期管理", font_size=14, color=IKB_BLUE)
add_text_box(slide, 0.8, 1.7, 11, 0.8, "ITSM客服端 · :5174", font_size=44, bold=True, color=DARK_BG)

# 截图网格
add_image(slide, "images/05-itsm-list.png", 0.8, 3.2, 3.7, 2.8)
add_image(slide, "images/06-itsm-accept.png", 4.8, 3.2, 3.7, 2.8)
add_image(slide, "images/07-itsm-detail.png", 8.8, 3.2, 3.7, 2.8)

add_text_box(slide, 0.8, 6.2, 11, 0.3, "四象限看板 · 状态流转 · SLA管理 · 实时通知 · 批量处理", font_size=12, color=GREY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 0.8, 7.1, 5, 0.3, "ITSM客服端", font_size=10, color=GREY)
add_text_box(slide, 10, 7.1, 2, 0.3, "08 / 12", font_size=10, color=GREY, alignment=PP_ALIGN.RIGHT)


# ==================== Slide 9: OPS统计 + 后台管理 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(0xF0, 0xF0, 0xEE))

add_text_box(slide, 0.8, 0.5, 3, 0.3, "08  ·  统计与管理", font_size=11, color=GREY)
add_text_box(slide, 0.8, 1.2, 5, 0.3, "数据分析 + 系统管理", font_size=14, color=IKB_BLUE)
add_text_box(slide, 0.8, 1.7, 11, 0.8, "OPS统计端 + 后台管理端", font_size=44, bold=True, color=DARK_BG)

# 左侧 OPS
add_text_box(slide, 0.8, 3.0, 5, 0.3, "OPS统计端 · :5176", font_size=14, color=IKB_BLUE)
add_image(slide, "images/08-ops-overview.png", 0.8, 3.5, 5.5, 3.2)

# 右侧 后台
add_text_box(slide, 7, 3.0, 5, 0.3, "后台管理端 · :5175", font_size=14, color=IKB_BLUE)
add_image(slide, "images/10-admin-users.png", 7, 3.5, 5.5, 3.2)

add_text_box(slide, 0.8, 7.1, 5, 0.3, "统计与管理", font_size=10, color=GREY)
add_text_box(slide, 10, 7.1, 2, 0.3, "09 / 12", font_size=10, color=GREY, alignment=PP_ALIGN.RIGHT)


# ==================== Slide 10: 关键数据 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, IKB_BLUE)

add_text_box(slide, 0.8, 0.5, 3, 0.3, "09  ·  关键数据", font_size=11, color=RGBColor(0x80, 0xA0, 0xD0))
add_text_box(slide, 0.8, 1.2, 5, 0.3, "项目成果量化", font_size=14, color=WHITE)
add_text_box(slide, 0.8, 1.7, 11, 0.8, "核心指标", font_size=44, bold=True, color=WHITE)

# 四个KPI
kpis = [
    ("75", "API端点", "RESTful接口"),
    ("73", "测试用例", "100%通过率"),
    ("4", "前端应用", "Vue 3独立应用"),
    ("100+", "知识库", "FAQ文档"),
]
for i, (value, label, note) in enumerate(kpis):
    x = 0.8 + i * 3.05
    add_shape_bg(slide, x, 3.2, 2.8, 2.8, RGBColor(0x00, 0x1A, 0x6A))
    add_text_box(slide, x + 0.3, 3.5, 2.2, 0.3, label, font_size=12, color=RGBColor(0x80, 0xA0, 0xD0))
    add_text_box(slide, x + 0.3, 4.0, 2.2, 1.2, value, font_size=64, bold=True, color=WHITE)
    add_text_box(slide, x + 0.3, 5.3, 2.2, 0.3, note, font_size=12, color=RGBColor(0x80, 0xA0, 0xD0))

add_text_box(slide, 0.8, 7.1, 5, 0.3, "关键数据", font_size=10, color=RGBColor(0x80, 0xA0, 0xD0))
add_text_box(slide, 10, 7.1, 2, 0.3, "10 / 12", font_size=10, color=RGBColor(0x80, 0xA0, 0xD0), alignment=PP_ALIGN.RIGHT)


# ==================== Slide 11: 性能优化 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GREY)

add_text_box(slide, 0.8, 0.5, 3, 0.3, "10  ·  性能优化", font_size=11, color=GREY)
add_text_box(slide, 0.8, 1.2, 5, 0.3, "系统优化措施", font_size=14, color=IKB_BLUE)
add_text_box(slide, 0.8, 1.7, 11, 0.8, "性能与稳定性", font_size=44, bold=True, color=DARK_BG)

# 优化时间线
opts = [
    ("数据库", "性能索引 · N+1修复 · 批量查询优化"),
    ("缓存", "Redis双写 · 分类缓存 · 权限缓存60s"),
    ("并发", "批量接口100张 · 限流300次/分 · 自动重试"),
    ("前端", "搜索防抖 · 智能滚动 · 429去重提示"),
    ("AI", "RAG去重 · 反幻觉过滤 · 流式输出"),
]
for i, (title, desc) in enumerate(opts):
    x = 0.8 + i * 2.45
    add_shape_bg(slide, x, 3.2, 2.2, 3.0, WHITE)
    add_shape_bg(slide, x, 3.2, 2.2, 0.06, IKB_BLUE)
    add_text_box(slide, x + 0.2, 3.5, 1.8, 0.3, title, font_size=16, bold=True, color=IKB_BLUE)
    add_text_box(slide, x + 0.2, 4.2, 1.8, 1.8, desc, font_size=12, color=RGBColor(0x52, 0x52, 0x52))

add_text_box(slide, 0.8, 7.1, 5, 0.3, "性能优化", font_size=10, color=GREY)
add_text_box(slide, 10, 7.1, 2, 0.3, "11 / 12", font_size=10, color=GREY, alignment=PP_ALIGN.RIGHT)


# ==================== Slide 12: 总结 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_text_box(slide, 0.8, 0.5, 3, 0.3, "11  ·  总结", font_size=11, color=GREY)
add_text_box(slide, 0.8, 1.2, 5, 0.3, "项目亮点", font_size=14, color=IKB_BLUE)
add_text_box(slide, 0.8, 2.0, 5, 2.0, "IT服务台\nAI赋能", font_size=60, bold=True, color=WHITE)
add_text_box(slide, 0.8, 4.5, 5, 1.5, "从传统工单系统升级为AI智能客服平台，实现问题自动解答、知识库沉淀、多轮精准对话", font_size=18, color=RGBColor(0xBB, 0xBB, 0xBB))

# 亮点卡片
highlights = [
    ("RAG", "检索增强生成"),
    ("Memory", "三层会话记忆"),
    ("SSE", "流式输出"),
    ("73/73", "测试全过"),
]
for i, (title, desc) in enumerate(highlights):
    row = i // 2
    col = i % 2
    x = 7 + col * 2.8
    y = 2.0 + row * 2.0
    add_shape_bg(slide, x, y, 2.5, 1.6, RGBColor(0x1A, 0x1A, 0x1A))
    add_text_box(slide, x + 0.2, y + 0.2, 2.1, 0.4, title, font_size=18, bold=True, color=IKB_BLUE)
    add_text_box(slide, x + 0.2, y + 0.8, 2.1, 0.5, desc, font_size=14, color=RGBColor(0x73, 0x73, 0x73))

add_text_box(slide, 0.8, 7.1, 5, 0.3, "感谢聆听", font_size=10, color=GREY)
add_text_box(slide, 10, 7.1, 2, 0.3, "12 / 12", font_size=10, color=GREY, alignment=PP_ALIGN.RIGHT)


# 保存
output_path = "C:\\work\\program_last\\项目详细\\ppt\\公司桌面IT服务台-项目报告.pptx"
prs.save(output_path)
print(f"PPT已生成: {output_path}")
